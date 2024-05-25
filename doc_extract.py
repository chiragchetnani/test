import PyPDF2
from docx import Document
import os 
from pptx import Presentation
import torch


#extracting from pdf
def pdf(path) : 

    pdf_file = open(path , 'rb')
    reader = PyPDF2.PdfReader(pdf_file)
    num_pages = len(reader.pages)
    text = ''

    for page_num in range(num_pages):
        page = reader.pages[page_num]
        text += page.extract_text()

    pdf_file.close()

    return text

#extracting from docx
def docx_(path) : 

    doc = Document(path)
    text = ''

    for paragraph in doc.paragraphs : text += paragraph.text

    return text

import docx

# #extracting from doc
# def doc(path) : 

#     doc = Document(path)
#     new_path = f'{path}x'
#     doc.save(new_path)

#     text = extract_text_from_docx(new_path)

#     os.remove(new_path)

#     return text

#extracting from pptx
def pptx(path) : 

    pr = Presentation(path)
    text = ''

    for slide in pr.slides : 

        for shape in slide.shapes : 
            if shape.has_text_frame : 
                for paragraph in shape.text_frame.paragraphs : 
                    for run in paragraph.runs : 
                      text += run.text

    return text

#extracting from csv
def csv(path) : 

    text = open(path).read()

#extracting from txt
def txt(path) : 

    text = open(path).read()

    return text
